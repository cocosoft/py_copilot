"""文件处理器服务"""
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Dict, List, Optional, Any


class FileProcessorInterface(ABC):
    """文件处理器接口"""
    
    @abstractmethod
    def can_process(self, file_path: Path, file_type: str) -> bool:
        """判断是否可以处理该文件"""
        pass
    
    @abstractmethod
    def process(self, file_path: Path, file_name: str) -> Dict[str, Any]:
        """处理文件并返回内容"""
        pass


class TextFileProcessor(FileProcessorInterface):
    """文本文件处理器"""
    
    def can_process(self, file_path: Path, file_type: str) -> bool:
        return file_type == 'text' or file_path.suffix in ['.txt', '.md', '.json', '.yaml', '.yml', '.csv']
    
    def process(self, file_path: Path, file_name: str) -> Dict[str, Any]:
        """处理文本文件"""
        content = ""
        error_message = None
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    content = f.read()
            except Exception as e:
                error_message = f"编码错误: {str(e)}"
        except Exception as e:
            error_message = f"读取失败: {str(e)}"
        
        if error_message:
            return {
                'filename': file_name,
                'content': f"[文件内容读取失败: {error_message}]",
                'type': 'text',
                'error': error_message
            }
        
        return {
            'filename': file_name,
            'content': content,
            'type': 'text'
        }


class WordFileProcessor(FileProcessorInterface):
    """Word文件处理器"""
    
    def can_process(self, file_path: Path, file_type: str) -> bool:
        return file_type == 'word' or file_path.suffix in ['.docx', '.doc']
    
    def process(self, file_path: Path, file_name: str) -> Dict[str, Any]:
        """处理Word文件"""
        content = ""
        error_message = None
        
        try:
            if file_path.suffix == '.docx':
                from docx import Document
                doc = Document(file_path)
                content = '\n'.join([para.text for para in doc.paragraphs])
            elif file_path.suffix == '.doc':
                content = "[.doc文件内容需要特殊处理]"
            else:
                content = f"[Word文件内容需要特殊处理: {file_path.suffix}]"
        except ImportError:
            error_message = "缺少python-docx库"
        except Exception as e:
            error_message = f"读取失败: {str(e)}"
        
        if error_message:
            return {
                'filename': file_name,
                'content': f"[Word文件内容读取失败: {error_message}]",
                'type': 'word',
                'error': error_message
            }
        
        return {
            'filename': file_name,
            'content': content,
            'type': 'word'
        }


class ExcelFileProcessor(FileProcessorInterface):
    """Excel文件处理器"""
    
    def can_process(self, file_path: Path, file_type: str) -> bool:
        return file_type == 'excel' or file_path.suffix in ['.xlsx', '.xls']
    
    def process(self, file_path: Path, file_name: str) -> Dict[str, Any]:
        """处理Excel文件"""
        content = ""
        error_message = None
        
        try:
            import pandas as pd
            excel_data = []
            xl_file = pd.ExcelFile(file_path)
            for sheet_name in xl_file.sheet_names:
                df = pd.read_excel(xl_file, sheet_name)
                excel_data.append(f"=== Sheet: {sheet_name} ===")
                excel_data.append(df.to_string())
            content = '\n'.join(excel_data)
        except ImportError:
            error_message = "缺少pandas库"
        except Exception as e:
            error_message = f"读取失败: {str(e)}"
        
        if error_message:
            return {
                'filename': file_name,
                'content': f"[Excel文件内容读取失败: {error_message}]",
                'type': 'excel',
                'error': error_message
            }
        
        return {
            'filename': file_name,
            'content': content,
            'type': 'excel'
        }


class PptFileProcessor(FileProcessorInterface):
    """PPT文件处理器"""
    
    def can_process(self, file_path: Path, file_type: str) -> bool:
        return file_type == 'ppt' or file_path.suffix in ['.pptx']
    
    def process(self, file_path: Path, file_name: str) -> Dict[str, Any]:
        """处理PPT文件"""
        content = ""
        error_message = None
        
        try:
            from pptx import Presentation
            ppt_data = []
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides):
                slide_content = []
                slide_content.append(f"=== Slide {i+1} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        slide_content.append(shape.text)
                
                if slide_content:
                    ppt_data.extend(slide_content)
                    ppt_data.append("")  # 添加空行分隔幻灯片
            
            content = '\n'.join(ppt_data)
        except ImportError:
            error_message = "缺少python-pptx库"
        except Exception as e:
            error_message = f"读取失败: {str(e)}"
        
        if error_message:
            return {
                'filename': file_name,
                'content': f"[PPT文件内容读取失败: {error_message}]",
                'type': 'ppt',
                'error': error_message
            }
        
        return {
            'filename': file_name,
            'content': content,
            'type': 'ppt'
        }


class PdfFileProcessor(FileProcessorInterface):
    """PDF文件处理器"""
    
    def can_process(self, file_path: Path, file_type: str) -> bool:
        return file_type == 'pdf' or file_path.suffix == '.pdf'
    
    def process(self, file_path: Path, file_name: str) -> Dict[str, Any]:
        """处理PDF文件"""
        content = ""
        error_message = None
        
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        content += page_text + '\n'
        except ImportError:
            error_message = "缺少pdfplumber库"
        except Exception as e:
            error_message = f"读取失败: {str(e)}"
        
        if error_message:
            return {
                'filename': file_name,
                'content': f"[PDF文件内容读取失败: {error_message}]",
                'type': 'pdf',
                'error': error_message
            }
        
        return {
            'filename': file_name,
            'content': content,
            'type': 'pdf'
        }


class DefaultFileProcessor(FileProcessorInterface):
    """默认文件处理器"""
    
    def can_process(self, file_path: Path, file_type: str) -> bool:
        return True  # 默认处理器可以处理任何文件
    
    def process(self, file_path: Path, file_name: str) -> Dict[str, Any]:
        """处理默认文件"""
        file_size = file_path.stat().st_size if file_path.exists() else 0
        return {
            'filename': file_name,
            'content': f"[文件类型: {file_path.suffix.lstrip('.')}, 大小: {file_size} 字节]",
            'type': 'other'
        }


class FileProcessorService:
    """文件处理器服务"""
    
    def __init__(self):
        """初始化文件处理器服务"""
        self.processors: List[FileProcessorInterface] = [
            TextFileProcessor(),
            WordFileProcessor(),
            ExcelFileProcessor(),
            PptFileProcessor(),
            PdfFileProcessor(),
            DefaultFileProcessor()
        ]
    
    def process_file(self, file_path: Path, file_name: str, file_type: str) -> Dict[str, Any]:
        """
        处理文件并返回内容
        
        Args:
            file_path: 文件路径
            file_name: 文件名
            file_type: 文件类型
            
        Returns:
            包含文件内容的字典
        """
        # 查找适合的处理器
        for processor in self.processors:
            if processor.can_process(file_path, file_type):
                print(f"使用 {processor.__class__.__name__} 处理文件: {file_name}")
                return processor.process(file_path, file_name)
        
        # 如果没有找到适合的处理器，使用默认处理器
        default_processor = DefaultFileProcessor()
        print(f"使用 DefaultFileProcessor 处理文件: {file_name}")
        return default_processor.process(file_path, file_name)
    
    def add_processor(self, processor: FileProcessorInterface):
        """
        添加自定义文件处理器
        
        Args:
            processor: 文件处理器实例
        """
        self.processors.insert(0, processor)  # 新处理器优先级更高
    
    def get_supported_file_types(self) -> List[str]:
        """
        获取支持的文件类型
        
        Returns:
            支持的文件类型列表
        """
        return ['text', 'word', 'excel', 'ppt', 'pdf', 'other']


# 创建全局文件处理器服务实例
file_processor_service = FileProcessorService()
