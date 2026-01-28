import PyPDF2
import docx
import openpyxl
import pptx
import pandas as pd
from PIL import Image
import io
import os
import json
import time
from pathlib import Path
from typing import Optional, Dict, List, Any
from dataclasses import dataclass
from enum import Enum


class FileType(Enum):
    """文件类型枚举"""
    PDF = "pdf"
    WORD = "word"
    EXCEL = "excel"
    POWERPOINT = "powerpoint"
    TEXT = "text"
    IMAGE = "image"
    WPS_WORD = "wps_word"
    WPS_EXCEL = "wps_excel"
    WPS_POWERPOINT = "wps_powerpoint"
    UNKNOWN = "unknown"


class ParseStatus(Enum):
    """解析状态枚举"""
    SUCCESS = "success"
    PARTIAL = "partial"
    FAILED = "failed"


@dataclass
class DocumentMetadata:
    """文档元数据"""
    file_name: str
    file_size: int
    file_type: FileType
    page_count: Optional[int] = None
    author: Optional[str] = None
    created_time: Optional[str] = None
    modified_time: Optional[str] = None
    language: Optional[str] = None


@dataclass
class ParseResult:
    """解析结果"""
    status: ParseStatus
    content: str
    metadata: DocumentMetadata
    processing_time: float
    error_message: Optional[str] = None
    extracted_entities: Optional[List[Dict[str, Any]]] = None
    confidence_score: Optional[float] = None


class DocumentParser:
    SUPPORTED_FORMATS = [
        '.pdf', '.docx', '.doc', '.txt', '.xlsx', '.xls', '.pptx', '.ppt', 
        '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.wps', '.et', '.dps',
        '.csv', '.json', '.xml', '.html', '.md', '.rtf', '.odt', '.ods', '.odp'
    ]
    
    # OCR支持的文件格式
    OCR_SUPPORTED_FORMATS = ['.pdf', '.png', '.jpg', '.jpeg', '.tiff', '.tif']
    
    # 大文件阈值（10MB）
    LARGE_FILE_THRESHOLD = 10 * 1024 * 1024
    
    @staticmethod
    def is_supported_format(file_path: str) -> bool:
        """检查文件格式是否支持"""
        return Path(file_path).suffix.lower() in DocumentParser.SUPPORTED_FORMATS
    
    @staticmethod
    def _is_ocr_supported(file_path: str) -> bool:
        """检查是否支持OCR"""
        return Path(file_path).suffix.lower() in DocumentParser.OCR_SUPPORTED_FORMATS
    
    @staticmethod
    def _is_large_file(file_path: str) -> bool:
        """检查是否为大文件"""
        try:
            file_size = os.path.getsize(file_path)
            return file_size > DocumentParser.LARGE_FILE_THRESHOLD
        except:
            return False
    
    @staticmethod
    def _extract_metadata(file_path: str, file_type: FileType) -> DocumentMetadata:
        """提取文档元数据"""
        file_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        
        # 获取文件创建和修改时间
        stat_info = os.stat(file_path)
        created_time = time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(stat_info.st_ctime))
        modified_time = time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(stat_info.st_mtime))
        
        return DocumentMetadata(
            file_name=file_name,
            file_size=file_size,
            file_type=file_type,
            created_time=created_time,
            modified_time=modified_time
        )
    
    @staticmethod
    def _extract_text_with_ocr(file_path: str) -> str:
        """使用OCR提取文本"""
        try:
            # 尝试导入OCR库
            import pytesseract
            from pdf2image import convert_from_path
            
            suffix = Path(file_path).suffix.lower()
            
            if suffix == '.pdf':
                # 处理PDF文件
                images = convert_from_path(file_path, dpi=300)
                text = ""
                for i, image in enumerate(images):
                    page_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                    text += f"=== 第{i+1}页 ===\n{page_text}\n\n"
                return text
            else:
                # 处理图像文件
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                return text
                
        except ImportError:
            return "OCR功能需要安装pytesseract和pdf2image库"
        except Exception as e:
            return f"OCR处理失败: {str(e)}"
    
    @staticmethod
    def parse_pdf(file_path: str, use_ocr: bool = False) -> str:
        """解析PDF文档"""
        start_time = time.time()
        
        try:
            # 首先尝试标准文本提取
            text = ""
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                
                # 提取文档信息
                info = reader.metadata
                author = info.get('/Author', '') if info else ''
                
                for page_num, page in enumerate(reader.pages, 1):
                    page_text = page.extract_text()
                    
                    # 检查文本提取是否有效
                    if len(page_text.strip()) < 50 and use_ocr:
                        # 文本提取效果不佳，尝试OCR
                        ocr_text = DocumentParser._extract_text_with_ocr(file_path)
                        if ocr_text and not ocr_text.startswith("OCR"):
                            text += f"=== 第{page_num}页 (OCR提取) ===\n{ocr_text}\n\n"
                        else:
                            text += f"=== 第{page_num}页 ===\n{page_text}\n\n"
                    else:
                        text += f"=== 第{page_num}页 ===\n{page_text}\n\n"
            
            # 如果标准提取失败且支持OCR，尝试OCR
            if not text.strip() and DocumentParser._is_ocr_supported(file_path):
                text = DocumentParser._extract_text_with_ocr(file_path)
            
            processing_time = time.time() - start_time
            
            if not text.strip():
                return f"PDF解析完成，但未提取到文本内容。处理时间: {processing_time:.2f}秒"
            
            return text
            
        except Exception as e:
            processing_time = time.time() - start_time
            
            # 尝试OCR作为备用方案
            if DocumentParser._is_ocr_supported(file_path):
                try:
                    ocr_text = DocumentParser._extract_text_with_ocr(file_path)
                    if ocr_text and not ocr_text.startswith("OCR"):
                        return f"标准解析失败，使用OCR提取:\n{ocr_text}"
                except:
                    pass
            
            raise ValueError(f"PDF解析失败: {str(e)}")
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """解析Word文档"""
        start_time = time.time()
        
        try:
            doc = docx.Document(file_path)
            text = ""
            
            # 提取文档属性
            core_props = doc.core_properties
            author = core_props.author or "未知作者"
            created_time = core_props.created.strftime('%Y-%m-%d %H:%M:%S') if core_props.created else "未知时间"
            
            text += f"=== 文档信息 ===\n"
            text += f"标题: {core_props.title or '无标题'}\n"
            text += f"作者: {author}\n"
            text += f"创建时间: {created_time}\n"
            text += f"页数: 约{len(doc.paragraphs)//50 + 1}页\n\n"
            
            # 提取正文内容
            for i, paragraph in enumerate(doc.paragraphs):
                if paragraph.text.strip():
                    # 检测标题样式
                    if paragraph.style.name.startswith('Heading'):
                        text += f"\n# {paragraph.text}\n"
                    else:
                        text += paragraph.text + "\n"
            
            # 提取表格内容
            if doc.tables:
                text += f"\n=== 表格内容 ===\n"
                for table_num, table in enumerate(doc.tables, 1):
                    text += f"表格{table_num}:\n"
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        text += f"{row_text}\n"
                    text += "\n"
            
            processing_time = time.time() - start_time
            
            if not text.strip():
                return f"Word文档解析完成，但未提取到文本内容。处理时间: {processing_time:.2f}秒"
            
            return text
            
        except Exception as e:
            processing_time = time.time() - start_time
            raise ValueError(f"Word文档解析失败: {str(e)}")
    
    @staticmethod
    def parse_txt(file_path: str) -> str:
        """解析文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    return file.read()
            except Exception as e:
                raise ValueError(f"文本文件解析失败: {str(e)}")
    
    @staticmethod
    def parse_excel(file_path: str) -> str:
        """解析Excel文件"""
        try:
            # 使用pandas读取Excel文件
            excel_file = pd.ExcelFile(file_path)
            text = ""
            
            # 遍历所有工作表
            for sheet_name in excel_file.sheet_names:
                df = excel_file.parse(sheet_name)
                text += f"=== 工作表: {sheet_name} ===\n"
                
                # 添加列名
                text += "\t".join([str(col) for col in df.columns]) + "\n"
                
                # 添加数据行
                for _, row in df.iterrows():
                    text += "\t".join([str(cell) for cell in row]) + "\n"
                
                text += "\n"
            
            return text
        except Exception as e:
            raise ValueError(f"Excel文件解析失败: {str(e)}")
    
    @staticmethod
    def parse_powerpoint(file_path: str) -> str:
        """解析PowerPoint文件"""
        try:
            prs = pptx.Presentation(file_path)
            text = ""
            
            # 遍历所有幻灯片
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"=== 幻灯片 {slide_num} ===\n"
                
                # 获取幻灯片中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text += shape.text + "\n"
                
                text += "\n"
            
            return text
        except Exception as e:
            raise ValueError(f"PowerPoint文件解析失败: {str(e)}")
    
    @staticmethod
    def parse_image(file_path: str) -> str:
        """解析图片文件"""
        try:
            with Image.open(file_path) as img:
                width, height = img.size
                format = img.format
                mode = img.mode
                
                # 返回图片基本信息
                return f"=== 图片信息 ===\n" \
                       f"文件名: {os.path.basename(file_path)}\n" \
                       f"格式: {format}\n" \
                       f"尺寸: {width} x {height} 像素\n" \
                       f"色彩模式: {mode}\n"
        except Exception as e:
            raise ValueError(f"图片文件解析失败: {str(e)}")

    @staticmethod
    def parse_wps(file_path: str) -> str:
        """解析WPS文件"""
        try:
            suffix = Path(file_path).suffix.lower()
            
            if suffix == '.wps':
                # WPS文字文档 - 尝试使用docx库解析
                try:
                    doc = docx.Document(file_path)
                    text = ""
                    for paragraph in doc.paragraphs:
                        text += paragraph.text + "\n"
                    return f"=== WPS文字文档内容 ===\n{text}"
                except Exception:
                    # 如果docx库无法解析，返回文件信息
                    file_size = os.path.getsize(file_path)
                    return f"=== WPS文字文档信息 ===\n" \
                           f"文件名: {os.path.basename(file_path)}\n" \
                           f"文件大小: {file_size} 字节\n" \
                           f"说明: 此WPS文件格式需要特殊处理\n"
            
            elif suffix == '.et':
                # WPS表格文档 - 尝试使用openpyxl库解析
                try:
                    workbook = openpyxl.load_workbook(file_path)
                    text = ""
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text += f"=== 工作表: {sheet_name} ===\n"
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                            text += row_text + "\n"
                        text += "\n"
                    return text
                except Exception:
                    # 如果openpyxl库无法解析，返回文件信息
                    file_size = os.path.getsize(file_path)
                    return f"=== WPS表格文档信息 ===\n" \
                           f"文件名: {os.path.basename(file_path)}\n" \
                           f"文件大小: {file_size} 字节\n" \
                           f"说明: 此WPS表格文件格式需要特殊处理\n"
            
            elif suffix == '.dps':
                # WPS演示文档 - 尝试使用pptx库解析
                try:
                    presentation = pptx.Presentation(file_path)
                    text = ""
                    for i, slide in enumerate(presentation.slides):
                        text += f"=== 幻灯片 {i+1} ===\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                text += shape.text + "\n"
                        text += "\n"
                    return text
                except Exception:
                    # 如果pptx库无法解析，返回文件信息
                    file_size = os.path.getsize(file_path)
                    return f"=== WPS演示文档信息 ===\n" \
                           f"文件名: {os.path.basename(file_path)}\n" \
                           f"文件大小: {file_size} 字节\n" \
                           f"说明: 此WPS演示文件格式需要特殊处理\n"
            
            else:
                raise ValueError(f"不支持的WPS文件格式: {suffix}")
                
        except Exception as e:
            raise ValueError(f"WPS文件解析失败: {str(e)}")
    
    @staticmethod
    def parse_document(file_path: str) -> str:
        """统一文档解析入口"""
        suffix = Path(file_path).suffix.lower()
        if suffix == '.pdf':
            return DocumentParser.parse_pdf(file_path)
        elif suffix in ['.docx', '.doc']:
            return DocumentParser.parse_docx(file_path)
        elif suffix == '.txt':
            return DocumentParser.parse_txt(file_path)
        elif suffix in ['.xlsx', '.xls']:
            return DocumentParser.parse_excel(file_path)
        elif suffix in ['.pptx', '.ppt']:
            return DocumentParser.parse_powerpoint(file_path)
        elif suffix in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            return DocumentParser.parse_image(file_path)
        elif suffix in ['.wps', '.et', '.dps']:
            return DocumentParser.parse_wps(file_path)
        else:
            raise ValueError(f"不支持的文档格式: {suffix}")
    
    @staticmethod
    def parse_document_advanced(file_path: str, use_ocr: bool = True, extract_entities: bool = False) -> ParseResult:
        """高级文档解析接口"""
        start_time = time.time()
        
        try:
            # 检测文件类型
            file_type = DocumentParser._detect_file_type(file_path)
            
            # 提取元数据
            metadata = DocumentParser._extract_metadata(file_path, file_type)
            
            # 根据文件类型选择解析方法
            content = ""
            status = ParseStatus.SUCCESS
            error_message = None
            
            try:
                if file_type == FileType.PDF:
                    content = DocumentParser.parse_pdf(file_path, use_ocr)
                elif file_type == FileType.WORD:
                    content = DocumentParser.parse_docx(file_path)
                elif file_type == FileType.EXCEL:
                    content = DocumentParser.parse_excel(file_path)
                elif file_type == FileType.POWERPOINT:
                    content = DocumentParser.parse_powerpoint(file_path)
                elif file_type == FileType.TEXT:
                    content = DocumentParser.parse_txt(file_path)
                elif file_type == FileType.IMAGE:
                    content = DocumentParser.parse_image(file_path)
                elif file_type in [FileType.WPS_WORD, FileType.WPS_EXCEL, FileType.WPS_POWERPOINT]:
                    content = DocumentParser.parse_wps(file_path)
                else:
                    content = f"不支持的文件类型: {file_type.value}"
                    status = ParseStatus.FAILED
            
            except Exception as e:
                content = f"解析失败: {str(e)}"
                status = ParseStatus.FAILED
                error_message = str(e)
            
            # 计算处理时间
            processing_time = time.time() - start_time
            
            # 计算置信度分数
            confidence_score = DocumentParser._calculate_confidence_score(content, file_type)
            
            # 提取实体（如果启用）
            extracted_entities = None
            if extract_entities and status == ParseStatus.SUCCESS:
                extracted_entities = DocumentParser._extract_entities(content)
            
            return ParseResult(
                status=status,
                content=content,
                metadata=metadata,
                processing_time=processing_time,
                error_message=error_message,
                extracted_entities=extracted_entities,
                confidence_score=confidence_score
            )
            
        except Exception as e:
            processing_time = time.time() - start_time
            
            # 创建基础元数据
            file_name = os.path.basename(file_path)
            file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
            metadata = DocumentMetadata(
                file_name=file_name,
                file_size=file_size,
                file_type=FileType.UNKNOWN
            )
            
            return ParseResult(
                status=ParseStatus.FAILED,
                content=f"解析过程出错: {str(e)}",
                metadata=metadata,
                processing_time=processing_time,
                error_message=str(e)
            )
    
    @staticmethod
    def _detect_file_type(file_path: str) -> FileType:
        """检测文件类型"""
        suffix = Path(file_path).suffix.lower()
        
        if suffix == '.pdf':
            return FileType.PDF
        elif suffix in ['.docx', '.doc']:
            return FileType.WORD
        elif suffix in ['.xlsx', '.xls']:
            return FileType.EXCEL
        elif suffix in ['.pptx', '.ppt']:
            return FileType.POWERPOINT
        elif suffix in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.rtf']:
            return FileType.TEXT
        elif suffix in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif']:
            return FileType.IMAGE
        elif suffix == '.wps':
            return FileType.WPS_WORD
        elif suffix == '.et':
            return FileType.WPS_EXCEL
        elif suffix == '.dps':
            return FileType.WPS_POWERPOINT
        else:
            return FileType.UNKNOWN
    
    @staticmethod
    def _calculate_confidence_score(content: str, file_type: FileType) -> float:
        """计算解析置信度分数"""
        if not content or len(content.strip()) < 10:
            return 0.0
        
        # 基于内容长度和质量计算分数
        content_length = len(content.strip())
        
        # 基础分数
        base_score = min(content_length / 1000, 1.0)
        
        # 基于文本质量调整分数
        # 计算中文字符比例（中文文档通常有更多有效内容）
        chinese_chars = sum(1 for char in content if '\u4e00' <= char <= '\u9fff')
        chinese_ratio = chinese_chars / len(content) if content else 0
        
        # 计算有效字符比例（非空格、标点）
        valid_chars = sum(1 for char in content if char.isalnum() or '\u4e00' <= char <= '\u9fff')
        valid_ratio = valid_chars / len(content) if content else 0
        
        # 综合分数
        confidence = base_score * 0.4 + chinese_ratio * 0.3 + valid_ratio * 0.3
        
        return min(confidence, 1.0)
    
    @staticmethod
    def _extract_entities(content: str) -> List[Dict[str, Any]]:
        """提取文本中的实体"""
        entities = []
        
        # 简单的中文实体提取（基于关键词）
        # 在实际应用中应该使用NLP库如spaCy或jieba
        
        # 提取人名（简单规则）
        import re
        
        # 中文人名模式（2-4个中文字符）
        chinese_name_pattern = r'[\u4e00-\u9fff]{2,4}'
        names = re.findall(chinese_name_pattern, content)
        
        for name in set(names):
            entities.append({
                'type': 'PERSON',
                'text': name,
                'start': content.find(name),
                'end': content.find(name) + len(name)
            })
        
        # 提取组织名（包含"公司"、"集团"等）
        org_pattern = r'[\u4e00-\u9fff]+(公司|集团|企业|银行|医院|学校)'
        orgs = re.findall(org_pattern, content)
        
        for org in set(orgs):
            entities.append({
                'type': 'ORGANIZATION',
                'text': org,
                'start': content.find(org),
                'end': content.find(org) + len(org)
            })
        
        # 提取地点（包含"省"、"市"、"区"等）
        location_pattern = r'[\u4e00-\u9fff]+(省|市|区|县|街道|路)'
        locations = re.findall(location_pattern, content)
        
        for location in set(locations):
            entities.append({
                'type': 'LOCATION',
                'text': location,
                'start': content.find(location),
                'end': content.find(location) + len(location)
            })
        
        return entities
    
    @staticmethod
    def batch_parse_documents(file_paths: List[str], use_ocr: bool = True, max_workers: int = 4) -> List[ParseResult]:
        """批量解析文档"""
        import concurrent.futures
        
        results = []
        
        # 使用线程池并行处理
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 提交所有任务
            future_to_file = {
                executor.submit(DocumentParser.parse_document_advanced, file_path, use_ocr, False): file_path 
                for file_path in file_paths
            }
            
            # 收集结果
            for future in concurrent.futures.as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    result = future.result()
                    results.append(result)
                except Exception as e:
                    # 创建错误结果
                    file_name = os.path.basename(file_path)
                    file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                    metadata = DocumentMetadata(
                        file_name=file_name,
                        file_size=file_size,
                        file_type=FileType.UNKNOWN
                    )
                    
                    error_result = ParseResult(
                        status=ParseStatus.FAILED,
                        content=f"批量处理失败: {str(e)}",
                        metadata=metadata,
                        processing_time=0.0,
                        error_message=str(e)
                    )
                    results.append(error_result)
        
        # 按处理时间排序
        results.sort(key=lambda x: x.processing_time)
        
        return results
    
    @staticmethod
    def parse_large_document(file_path: str, chunk_size: int = 10, use_ocr: bool = True) -> List[ParseResult]:
        """解析大文档（分块处理）"""
        if not DocumentParser._is_large_file(file_path):
            # 如果不是大文件，直接解析
            return [DocumentParser.parse_document_advanced(file_path, use_ocr)]
        
        file_type = DocumentParser._detect_file_type(file_path)
        
        if file_type != FileType.PDF:
            # 目前只支持PDF大文件分块处理
            return [DocumentParser.parse_document_advanced(file_path, use_ocr)]
        
        results = []
        
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                total_pages = len(reader.pages)
                
                # 分块处理
                for start_page in range(0, total_pages, chunk_size):
                    end_page = min(start_page + chunk_size, total_pages)
                    
                    # 创建临时文件保存当前块
                    temp_file = f"{file_path}_chunk_{start_page}_{end_page}.pdf"
                    
                    try:
                        writer = PyPDF2.PdfWriter()
                        for page_num in range(start_page, end_page):
                            writer.add_page(reader.pages[page_num])
                        
                        with open(temp_file, 'wb') as output:
                            writer.write(output)
                        
                        # 解析当前块
                        chunk_result = DocumentParser.parse_document_advanced(temp_file, use_ocr)
                        
                        # 添加分块信息
                        chunk_result.metadata.page_count = end_page - start_page
                        chunk_result.content = f"=== 分块 {start_page+1}-{end_page} ===\n{chunk_result.content}"
                        
                        results.append(chunk_result)
                        
                    finally:
                        # 清理临时文件
                        if os.path.exists(temp_file):
                            os.remove(temp_file)
            
            return results
            
        except Exception as e:
            # 如果分块处理失败，回退到整体处理
            return [DocumentParser.parse_document_advanced(file_path, use_ocr)]
    
    @staticmethod
    def get_parsing_statistics(results: List[ParseResult]) -> Dict[str, Any]:
        """获取解析统计信息"""
        if not results:
            return {}
        
        total_files = len(results)
        successful_files = sum(1 for r in results if r.status == ParseStatus.SUCCESS)
        failed_files = total_files - successful_files
        
        total_processing_time = sum(r.processing_time for r in results)
        avg_processing_time = total_processing_time / total_files if total_files > 0 else 0
        
        total_content_length = sum(len(r.content) for r in results)
        avg_content_length = total_content_length / total_files if total_files > 0 else 0
        
        avg_confidence = sum(r.confidence_score or 0 for r in results) / total_files if total_files > 0 else 0
        
        return {
            'total_files': total_files,
            'successful_files': successful_files,
            'failed_files': failed_files,
            'success_rate': successful_files / total_files * 100 if total_files > 0 else 0,
            'total_processing_time': total_processing_time,
            'avg_processing_time': avg_processing_time,
            'total_content_length': total_content_length,
            'avg_content_length': avg_content_length,
            'avg_confidence': avg_confidence
        }
    
    @staticmethod
    def export_parse_results(results: List[ParseResult], output_format: str = 'json') -> str:
        """导出解析结果"""
        if output_format == 'json':
            # 转换为可序列化的字典
            export_data = []
            for result in results:
                result_dict = {
                    'status': result.status.value,
                    'content': result.content,
                    'metadata': {
                        'file_name': result.metadata.file_name,
                        'file_size': result.metadata.file_size,
                        'file_type': result.metadata.file_type.value,
                        'page_count': result.metadata.page_count,
                        'author': result.metadata.author,
                        'created_time': result.metadata.created_time,
                        'modified_time': result.metadata.modified_time,
                        'language': result.metadata.language
                    },
                    'processing_time': result.processing_time,
                    'error_message': result.error_message,
                    'extracted_entities': result.extracted_entities,
                    'confidence_score': result.confidence_score
                }
                export_data.append(result_dict)
            
            return json.dumps(export_data, ensure_ascii=False, indent=2)
        
        elif output_format == 'text':
            text_output = ""
            for i, result in enumerate(results, 1):
                text_output += f"=== 文件 {i}: {result.metadata.file_name} ===\n"
                text_output += f"状态: {result.status.value}\n"
                text_output += f"处理时间: {result.processing_time:.2f}秒\n"
                text_output += f"置信度: {result.confidence_score or 0:.2f}\n"
                
                if result.error_message:
                    text_output += f"错误信息: {result.error_message}\n"
                
                text_output += f"内容预览: {result.content[:200]}...\n\n"
            
            return text_output
        
        else:
            raise ValueError(f"不支持的导出格式: {output_format}")